#!/usr/bin/env python3
"""build_pptx.py — 把截图目录组装成 16:9 PPTX（每页一张满屏图）。
可选：把原始 HTML 作为 OLE 对象嵌入第一页右下角，双击即可打开。

用法: python3 build_pptx.py <imgDir> <output.pptx> [embed.html]
"""
import sys
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

img_dir = Path(sys.argv[1])
out_path = Path(sys.argv[2])
embed_html = Path(sys.argv[3]) if len(sys.argv) > 3 and sys.argv[3] else None

manifest = json.loads((img_dir / "manifest.json").read_text())
slides = manifest["slides"]

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]        # 完全空白版式

slide_objs = []
for name in slides:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(img_dir / name), 0, 0,
        width=prs.slide_width, height=prs.slide_height,
    )
    slide_objs.append(slide)


def make_icon(png_path: Path, label: str):
    """生成一个文档样式图标（带中文标签），供 OLE 对象显示用。"""
    from PIL import Image, ImageDraw, ImageFont
    W, H = 320, 380
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    # 文档主体（白底 + 折角）
    fold = 78
    body = [(24, 20), (W - 24 - fold, 20), (W - 24, 20 + fold), (W - 24, H - 78), (24, H - 78)]
    d.polygon(body, fill=(255, 255, 255, 255), outline=(210, 210, 210, 255))
    d.polygon([(W - 24 - fold, 20), (W - 24, 20 + fold), (W - 24 - fold, 20 + fold)],
              fill=(232, 232, 232, 255), outline=(210, 210, 210, 255))
    # 顶部品牌红条
    d.rectangle([24, 20, W - 24 - fold, 20 + 40], fill=(197, 32, 42, 255))

    def font(sz):
        for fp in ("/System/Library/Fonts/PingFang.ttc",
                   "/System/Library/Fonts/STHeiti Medium.ttc",
                   "/System/Library/Fonts/Supplemental/Songti.ttc"):
            try:
                return ImageFont.truetype(fp, sz)
            except Exception:
                continue
        return ImageFont.load_default()

    # 文档中央 "HTML"
    f_big = font(72)
    d.text((W / 2, 175), "HTML", font=f_big, fill=(197, 32, 42, 255), anchor="mm")
    # 几条示意文本线
    for i, y in enumerate((250, 282, 314)):
        d.line([(56, y), (W - 56 - (30 if i == 2 else 0), y)], fill=(200, 200, 200, 255), width=6)
    # 底部标签
    f_lbl = font(34)
    d.text((W / 2, H - 40), label, font=f_lbl, fill=(60, 60, 60, 255), anchor="mm")
    img.save(png_path)


if embed_html and embed_html.exists():
    from ole_package import build_ole_package
    icon_png = img_dir / "_ole_icon.png"
    make_icon(icon_png, "双击打开原件")
    # 把 HTML 封装成标准 OLE Package（Windows PowerPoint 双击可释放并用浏览器打开）
    # 文件名去掉空格/括号，降低 Windows 临时路径 shell 启动出错的风险
    import re
    safe = embed_html.stem
    safe = re.sub(r"[\s\(\)\[\]\{\}]+", "", safe) or "courseware"
    label = safe + ".html"
    pkg_bin = img_dir / "_ole_object.bin"
    pkg_bin.write_bytes(build_ole_package(embed_html.read_bytes(), label))
    iw, ih = Inches(1.15), Inches(1.35)
    left = prs.slide_width - iw - Inches(0.25)
    top = prs.slide_height - ih - Inches(0.2)
    slide_objs[0].shapes.add_ole_object(
        str(pkg_bin), prog_id="Package",
        left=left, top=top, width=iw, height=ih,
        icon_file=str(icon_png), icon_width=iw, icon_height=ih,
    )
    print(f"已在第 1 页嵌入原始 HTML (OLE Package): {label}")

out_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out_path))
print(f"已生成 {out_path}  共 {len(slides)} 页")
