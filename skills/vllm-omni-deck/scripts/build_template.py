#!/usr/bin/env python3
"""Build the canonical editable vLLM-Omni PowerPoint template."""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
)
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.util import Inches, Pt


FONT_NAME = "Arial"
ALLOWED_FONT_SIZES = frozenset({12, 18, 28, 36})
SLIDE_WIDTH_IN = 10.0
SLIDE_HEIGHT_IN = 5.625
SKILL_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_LOGO = SKILL_ROOT / "assets" / "vllm-omni-logo.png"


class TemplateValidationError(ValueError):
    """Raised when the generated template violates its design contract."""


@dataclass(frozen=True)
class Box:
    """A slide-relative rectangle measured in inches."""

    x: float
    y: float
    width: float
    height: float


@dataclass(frozen=True)
class Palette:
    """Fixed vLLM-Omni presentation colors."""

    navy: RGBColor = RGBColor(0x0D, 0x23, 0x40)
    blue: RGBColor = RGBColor(0x3B, 0x82, 0xF6)
    orange: RGBColor = RGBColor(0xF5, 0xA6, 0x23)
    cyan: RGBColor = RGBColor(0x45, 0xB7, 0xE8)
    purple: RGBColor = RGBColor(0x7C, 0x4D, 0xCC)
    ink: RGBColor = RGBColor(0x17, 0x20, 0x33)
    muted: RGBColor = RGBColor(0x5E, 0x6B, 0x7A)
    rule: RGBColor = RGBColor(0xD8, 0xE2, 0xEC)
    panel: RGBColor = RGBColor(0xF4, 0xF7, 0xFB)
    white: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)


@dataclass(frozen=True)
class CliArgs:
    """Command-line arguments."""

    output: Path
    force: bool


COLORS = Palette()


def add_text(
    slide: Slide,
    box: Box,
    text: str,
    font_size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.0,
) -> Shape:
    """Add a consistently styled text box."""

    if font_size not in ALLOWED_FONT_SIZES:
        raise TemplateValidationError(
            f"Unsupported font size {font_size}; use {sorted(ALLOWED_FONT_SIZES)}"
        )
    shape = slide.shapes.add_textbox(
        Inches(box.x),
        Inches(box.y),
        Inches(box.width),
        Inches(box.height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical_anchor
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_panel(
    slide: Slide,
    box: Box,
    fill_color: RGBColor,
    *,
    line_color: RGBColor | None = None,
    rounded: bool = True,
    line_width: float = 1.0,
) -> Shape:
    """Add a flat panel with an optional quiet outline."""

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(box.x),
        Inches(box.y),
        Inches(box.width),
        Inches(box.height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def add_rule(slide: Slide, box: Box, color: RGBColor) -> None:
    """Add a thin rectangular rule."""

    add_panel(slide, box, color, rounded=False)


def add_chevron(
    slide: Slide,
    box: Box,
    color: RGBColor,
) -> None:
    """Add a small flat chevron between flow stages."""

    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(box.x),
        Inches(box.y),
        Inches(box.width),
        Inches(box.height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def set_background(slide: Slide, color: RGBColor) -> None:
    """Set a solid slide background."""

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_brand_plate(slide: Slide, logo_path: Path) -> None:
    """Place the dark-slide wordmark inside a white plate."""

    add_panel(
        slide,
        Box(0.52, 0.34, 1.92, 0.50),
        COLORS.white,
        rounded=True,
    )
    slide.shapes.add_picture(
        str(logo_path),
        Inches(0.66),
        Inches(0.45),
        width=Inches(1.64),
    )


def add_dark_page_number(slide: Slide, page_number: int) -> None:
    """Add a page number to a dark slide."""

    add_text(
        slide,
        Box(9.30, 5.20, 0.24, 0.20),
        f"{page_number:02d}",
        12,
        COLORS.white,
        alignment=PP_ALIGN.RIGHT,
    )


def add_content_header(slide: Slide, kicker: str, title: str) -> None:
    """Add the shared content-slide title treatment."""

    add_text(slide, Box(0.52, 0.24, 4.60, 0.18), kicker, 12, COLORS.blue, bold=True)
    add_rule(slide, Box(0.52, 0.49, 0.06, 0.43), COLORS.orange)
    add_text(slide, Box(0.68, 0.45, 8.82, 0.50), title, 28, COLORS.ink, bold=True)


def add_content_footer(
    slide: Slide,
    logo_path: Path,
    page_number: int,
    source_text: str,
) -> None:
    """Add the content-slide source line, logo, and page number."""

    add_rule(slide, Box(0.50, 5.12, 9.00, 0.012), COLORS.rule)
    add_text(slide, Box(0.52, 5.23, 6.90, 0.20), source_text, 12, COLORS.muted)
    slide.shapes.add_picture(
        str(logo_path),
        Inches(8.03),
        Inches(5.25),
        width=Inches(1.18),
    )
    add_text(
        slide,
        Box(9.34, 5.23, 0.20, 0.20),
        f"{page_number:02d}",
        12,
        COLORS.muted,
        alignment=PP_ALIGN.RIGHT,
    )


def add_card(
    slide: Slide,
    box: Box,
    label: str,
    title: str,
    body: str,
    accent: RGBColor,
) -> None:
    """Add an editable content card."""

    add_panel(slide, box, COLORS.panel, line_color=COLORS.rule)
    add_rule(slide, Box(box.x, box.y, 0.07, box.height), accent)
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.15, box.width - 0.34, 0.18),
        label,
        12,
        accent,
        bold=True,
    )
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.44, box.width - 0.34, 0.30),
        title,
        18,
        COLORS.ink,
        bold=True,
    )
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.86, box.width - 0.34, box.height - 1.00),
        body,
        12,
        COLORS.muted,
    )


def add_cover(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build layout 1: cover / closing."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.navy)
    add_brand_plate(slide, logo_path)
    add_text(
        slide,
        Box(0.56, 1.22, 5.70, 1.30),
        "vLLM-Omni\nTechnical Briefing",
        36,
        COLORS.white,
        bold=True,
    )
    add_text(
        slide,
        Box(0.58, 2.78, 5.45, 0.48),
        "A reusable visual system for multimodal serving stories",
        18,
        COLORS.white,
    )
    add_rule(slide, Box(0.58, 3.47, 0.72, 0.07), COLORS.orange)
    add_text(
        slide,
        Box(0.58, 3.73, 4.70, 0.38),
        "Presenter / team  ·  Event or review  ·  Month YYYY",
        12,
        COLORS.white,
    )

    motif_x = 6.68
    motif_y = 1.38
    motif_items = (
        ("TEXT", COLORS.blue),
        ("IMAGE", COLORS.cyan),
        ("AUDIO", COLORS.orange),
        ("VIDEO", COLORS.purple),
    )
    for index, (label, color) in enumerate(motif_items):
        y = motif_y + index * 0.68
        add_panel(
            slide,
            Box(motif_x, y, 1.34, 0.46),
            color,
            rounded=True,
        )
        add_text(
            slide,
            Box(motif_x, y + 0.13, 1.34, 0.18),
            label,
            12,
            COLORS.white,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_chevron(slide, Box(8.16, y + 0.13, 0.30, 0.19), COLORS.rule)
    add_panel(
        slide,
        Box(8.55, 2.12, 0.90, 1.12),
        COLORS.white,
        rounded=True,
    )
    add_text(
        slide,
        Box(8.55, 2.49, 0.90, 0.28),
        "OMNI",
        18,
        COLORS.navy,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_rule(slide, Box(0.58, 5.04, 8.86, 0.012), COLORS.blue)
    add_text(
        slide,
        Box(0.58, 5.16, 5.90, 0.20),
        "POWERPOINT-FIRST  ·  EDITABLE TEMPLATE",
        12,
        COLORS.white,
        bold=True,
    )


def add_section(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build layout 2: section divider / key message."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.navy)
    add_brand_plate(slide, logo_path)
    add_text(
        slide,
        Box(0.58, 1.32, 4.90, 0.20),
        "01  /  ARCHITECTURE",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(0.58, 1.72, 5.10, 1.22),
        "One serving path.\nMany modalities.",
        36,
        COLORS.white,
        bold=True,
    )
    add_text(
        slide,
        Box(0.60, 3.20, 4.85, 0.74),
        "Use a section slide to reset attention and state the chapter's single "
        "most important idea.",
        18,
        COLORS.white,
    )

    flow_items = (
        ("INPUT", "TEXT\nIMAGE\nA/V", COLORS.blue),
        ("ENGINE", "ROUTE\nRUN\nSTREAM", COLORS.orange),
        ("OUTPUT", "TOKENS\nMEDIA\nEVENTS", COLORS.purple),
    )
    for index, (label, body, color) in enumerate(flow_items):
        x = 6.10 + index * 1.20
        add_panel(
            slide,
            Box(x, 1.72, 0.98, 1.78),
            COLORS.white,
            line_color=color,
        )
        add_rule(slide, Box(x, 1.72, 0.98, 0.08), color)
        add_text(
            slide,
            Box(x + 0.10, 2.02, 0.78, 0.18),
            label,
            12,
            color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            Box(x + 0.08, 2.33, 0.82, 0.82),
            body,
            12,
            COLORS.ink,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        if index < len(flow_items) - 1:
            add_chevron(slide, Box(x + 1.02, 2.46, 0.16, 0.24), COLORS.rule)
    add_dark_page_number(slide, 2)


def add_agenda(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build layout 3: agenda / summary."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide, "LAYOUT 03  ·  AGENDA / SUMMARY", "What this deck should answer"
    )
    cards = (
        (
            Box(0.55, 1.24, 4.28, 1.52),
            "01  /  EXPERIENCE",
            "What should feel simple?",
            "Name the user problem before describing the system.",
            COLORS.blue,
        ),
        (
            Box(5.07, 1.24, 4.38, 1.52),
            "02  /  ARCHITECTURE",
            "Where does work happen?",
            "Expose the few stage boundaries that explain the design.",
            COLORS.cyan,
        ),
        (
            Box(0.55, 3.00, 4.28, 1.52),
            "03  /  EVIDENCE",
            "What proves the claim?",
            "Use one comparison with explicit metrics and provenance.",
            COLORS.orange,
        ),
        (
            Box(5.07, 3.00, 4.38, 1.52),
            "04  /  ROADMAP",
            "What decision comes next?",
            "Sequence actions by dependency, ownership, and risk.",
            COLORS.purple,
        ),
    )
    for box, label, title, body, accent in cards:
        add_card(slide, box, label, title, body, accent)
    add_content_footer(
        slide,
        logo_path,
        3,
        "Template guidance  ·  Replace with the approved agenda",
    )


def add_general_content(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build layout 4: general content."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide, "LAYOUT 04  ·  GENERAL CONTENT", "One interface, many modalities"
    )

    add_panel(slide, Box(0.55, 1.23, 3.25, 3.58), COLORS.panel, line_color=COLORS.rule)
    add_text(
        slide,
        Box(0.78, 1.49, 2.78, 0.18),
        "KEY TAKEAWAY",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(0.78, 1.85, 2.70, 0.86),
        "Keep the user contract stable while execution stages evolve.",
        18,
        COLORS.ink,
        bold=True,
    )
    principles = (
        ("01", "One request model", COLORS.blue),
        ("02", "Explicit stage contracts", COLORS.cyan),
        ("03", "Observable outputs", COLORS.purple),
    )
    for index, (number, label, color) in enumerate(principles):
        y = 3.00 + index * 0.51
        add_panel(slide, Box(0.78, y, 0.42, 0.34), color)
        add_text(
            slide,
            Box(0.78, y + 0.08, 0.42, 0.18),
            number,
            12,
            COLORS.white,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text(
            slide, Box(1.36, y + 0.05, 2.12, 0.22), label, 12, COLORS.ink, bold=True
        )

    modality_items = (
        ("TEXT", COLORS.blue),
        ("IMAGE", COLORS.cyan),
        ("AUDIO", COLORS.orange),
        ("VIDEO", COLORS.purple),
    )
    for index, (label, color) in enumerate(modality_items):
        y = 1.38 + index * 0.70
        add_panel(slide, Box(4.15, y, 1.14, 0.43), COLORS.white, line_color=color)
        add_text(
            slide,
            Box(4.15, y + 0.12, 1.14, 0.18),
            label,
            12,
            color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_chevron(slide, Box(5.43, y + 0.11, 0.26, 0.20), COLORS.rule)

    add_panel(slide, Box(5.85, 1.38, 2.10, 2.55), COLORS.navy)
    add_text(
        slide,
        Box(6.09, 1.72, 1.62, 0.18),
        "vLLM-OMNI",
        12,
        COLORS.orange,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Box(6.05, 2.18, 1.70, 0.70),
        "Unified\nserving core",
        18,
        COLORS.white,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        Box(6.08, 3.23, 1.64, 0.22),
        "route · schedule · stream",
        12,
        COLORS.white,
        alignment=PP_ALIGN.CENTER,
    )
    add_chevron(slide, Box(8.10, 2.48, 0.30, 0.24), COLORS.rule)
    add_panel(slide, Box(8.55, 2.06, 0.90, 1.08), COLORS.panel, line_color=COLORS.rule)
    add_text(
        slide,
        Box(8.55, 2.39, 0.90, 0.40),
        "EDITABLE\nOUTPUTS",
        12,
        COLORS.ink,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Box(4.16, 4.33, 5.22, 0.24),
        "Use the visual region for relationships, not decorative imagery.",
        12,
        COLORS.muted,
    )
    add_content_footer(
        slide,
        logo_path,
        4,
        "Illustrative architecture  ·  Replace with approved content",
    )


def add_architecture(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build layout 5: architecture / technical diagram."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "LAYOUT 05  ·  ARCHITECTURE / DIAGRAM",
        "Stage contracts keep the pipeline composable",
    )
    add_text(
        slide,
        Box(0.55, 1.12, 8.90, 0.30),
        "Stable contracts let each stage evolve without rewriting the entire serving path.",
        18,
        COLORS.ink,
    )

    stages = (
        ("01  /  ROUTE", "Router", "validate · classify", COLORS.blue),
        ("02  /  ENCODE", "Adapters", "normalize · batch", COLORS.cyan),
        ("03  /  EXECUTE", "Engine", "schedule · cache", COLORS.orange),
        ("04  /  EMIT", "Outputs", "decode · stream", COLORS.purple),
    )
    stage_width = 1.90
    start_x = 0.55
    gap = 0.47
    for index, (label, title, detail, color) in enumerate(stages):
        x = start_x + index * (stage_width + gap)
        add_panel(
            slide, Box(x, 1.72, stage_width, 1.42), COLORS.panel, line_color=color
        )
        add_rule(slide, Box(x, 1.72, stage_width, 0.09), color)
        add_text(
            slide,
            Box(x + 0.16, 1.98, stage_width - 0.32, 0.18),
            label,
            12,
            color,
            bold=True,
        )
        add_text(
            slide,
            Box(x + 0.16, 2.30, stage_width - 0.32, 0.36),
            title,
            18,
            COLORS.ink,
            bold=True,
        )
        add_text(
            slide,
            Box(x + 0.16, 2.78, stage_width - 0.32, 0.20),
            detail,
            12,
            COLORS.muted,
        )
        if index < len(stages) - 1:
            add_chevron(
                slide,
                Box(x + stage_width + 0.10, 2.30, 0.25, 0.30),
                COLORS.rule,
            )

    add_panel(slide, Box(0.55, 3.54, 8.90, 1.10), COLORS.navy)
    add_text(
        slide,
        Box(0.79, 3.78, 1.65, 0.18),
        "CROSS-CUTTING",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(0.79, 4.12, 2.00, 0.26),
        "Stage contracts",
        18,
        COLORS.white,
        bold=True,
    )
    cross_cutting = (
        ("OBSERVABILITY", 3.30, COLORS.blue),
        ("BACKPRESSURE", 5.26, COLORS.cyan),
        ("ERROR BOUNDS", 7.20, COLORS.purple),
    )
    for label, x, color in cross_cutting:
        add_panel(slide, Box(x, 3.90, 1.70, 0.43), color)
        add_text(
            slide,
            Box(x, 4.02, 1.70, 0.18),
            label,
            12,
            COLORS.white,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
    add_content_footer(
        slide,
        logo_path,
        5,
        "Illustrative system structure  ·  Replace with approved architecture",
    )


def add_evidence_chart(slide: Slide) -> None:
    """Add a native editable comparison chart."""

    chart_data = ChartData()
    chart_data.categories = ("Text", "Image", "Audio")
    chart_data.add_series("Baseline", (1.0, 1.0, 1.0))
    chart_data.add_series("Optimized", (1.5, 1.4, 1.3))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.55),
        Inches(1.62),
        Inches(5.70),
        Inches(2.92),
        chart_data,
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.font.name = FONT_NAME
    chart.font.size = Pt(12)
    chart.category_axis.tick_labels.font.name = FONT_NAME
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.tick_labels.font.name = FONT_NAME
    chart.value_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 1.8
    chart.value_axis.major_unit = 0.5
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = COLORS.rule
    chart.plots[0].gap_width = 58
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.show_value = True
    chart.plots[0].data_labels.font.name = FONT_NAME
    chart.plots[0].data_labels.font.size = Pt(12)
    series_colors = (COLORS.muted, COLORS.blue)
    for series, color in zip(chart.series, series_colors, strict=True):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.format.line.fill.background()


def add_evidence(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build layout 6: evidence / benchmark / comparison."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide, "LAYOUT 06  ·  EVIDENCE", "Evidence should make the comparison obvious"
    )
    add_text(
        slide,
        Box(0.55, 1.18, 4.10, 0.18),
        "RELATIVE THROUGHPUT  ·  HIGHER IS BETTER",
        12,
        COLORS.muted,
        bold=True,
    )
    add_evidence_chart(slide)

    add_panel(slide, Box(6.56, 1.34, 2.89, 1.48), COLORS.panel, line_color=COLORS.rule)
    add_text(
        slide,
        Box(6.80, 1.56, 2.36, 0.22),
        "Read the claim first",
        18,
        COLORS.ink,
        bold=True,
    )
    add_text(
        slide,
        Box(6.80, 2.02, 2.30, 0.48),
        "Highlight one comparison and keep the baseline visible.",
        12,
        COLORS.muted,
    )
    add_panel(slide, Box(6.56, 3.03, 2.89, 1.48), COLORS.navy)
    add_text(slide, Box(6.80, 3.26, 1.20, 0.38), "1.4×", 28, COLORS.white, bold=True)
    add_text(
        slide,
        Box(6.80, 3.78, 2.30, 0.42),
        "Example conclusion\nReplace with approved evidence",
        12,
        COLORS.white,
    )
    add_panel(slide, Box(0.55, 4.68, 5.70, 0.28), COLORS.orange)
    add_text(
        slide,
        Box(0.55, 4.74, 5.70, 0.18),
        "ILLUSTRATIVE DATA  ·  NOT BENCHMARK RESULTS",
        12,
        COLORS.white,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_content_footer(
        slide,
        logo_path,
        6,
        "Illustrative data — not benchmark results",
    )


def add_roadmap_card(
    slide: Slide,
    box: Box,
    label: str,
    title: str,
    items: tuple[str, str],
    accent: RGBColor,
) -> None:
    """Add a roadmap phase card."""

    add_panel(slide, box, COLORS.panel, line_color=COLORS.rule)
    add_rule(slide, Box(box.x, box.y, box.width, 0.09), accent)
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.25, box.width - 0.40, 0.18),
        label,
        12,
        accent,
        bold=True,
    )
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.57, box.width - 0.40, 0.62),
        title,
        18,
        COLORS.ink,
        bold=True,
    )
    for index, item in enumerate(items):
        y = box.y + 1.45 + index * 0.66
        add_panel(slide, Box(box.x + 0.20, y, 0.30, 0.30), accent)
        add_text(
            slide,
            Box(box.x + 0.20, y + 0.06, 0.30, 0.18),
            f"{index + 1:02d}",
            12,
            COLORS.white,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            Box(box.x + 0.64, y + 0.02, box.width - 0.88, 0.34),
            item,
            12,
            COLORS.ink,
        )


def add_roadmap(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build layout 7: roadmap / status."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "LAYOUT 07  ·  ROADMAP / STATUS",
        "Sequence work by dependency, not excitement",
    )
    add_text(
        slide,
        Box(0.55, 1.12, 8.90, 0.28),
        "Use three horizons to make ownership, prerequisites, and decision gates visible.",
        18,
        COLORS.ink,
    )
    roadmap_cards = (
        (
            Box(0.55, 1.63, 2.80, 2.65),
            "NOW  /  VALIDATE",
            "Stabilize the contract",
            ("Confirm stage boundaries", "Instrument the critical path"),
            COLORS.blue,
        ),
        (
            Box(3.60, 1.63, 2.80, 2.65),
            "NEXT  /  SCALE",
            "Prove the bottleneck",
            ("Run controlled evidence", "Assign optimization owner"),
            COLORS.orange,
        ),
        (
            Box(6.65, 1.63, 2.80, 2.65),
            "LATER  /  EXPAND",
            "Broaden support",
            ("Add the next modality", "Retire temporary paths"),
            COLORS.purple,
        ),
    )
    for box, label, title, items, accent in roadmap_cards:
        add_roadmap_card(slide, box, label, title, items, accent)
    add_panel(slide, Box(0.55, 4.51, 8.90, 0.46), COLORS.navy)
    add_text(
        slide,
        Box(0.76, 4.64, 1.35, 0.18),
        "DECISION GATE",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(2.22, 4.60, 6.88, 0.24),
        "Advance only when evidence and ownership are explicit.",
        18,
        COLORS.white,
        bold=True,
    )
    add_content_footer(
        slide,
        logo_path,
        7,
        "Illustrative roadmap  ·  Replace with approved owners and dates",
    )


def build_presentation(logo_path: Path) -> PresentationType:
    """Build the complete seven-slide presentation in memory."""

    if not logo_path.is_file():
        raise FileNotFoundError(f"Missing vLLM-Omni logo: {logo_path}")
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    presentation.core_properties.title = "vLLM-Omni editable presentation template"
    presentation.core_properties.subject = "Seven reusable vLLM-Omni slide layouts"
    presentation.core_properties.author = "vLLM-Omni"
    add_cover(presentation, logo_path)
    add_section(presentation, logo_path)
    add_agenda(presentation, logo_path)
    add_general_content(presentation, logo_path)
    add_architecture(presentation, logo_path)
    add_evidence(presentation, logo_path)
    add_roadmap(presentation, logo_path)
    return presentation


def slide_text(slide: Slide) -> str:
    """Return visible shape text for structural validation."""

    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text)
    return "\n".join(parts)


def validate_template(output_path: Path) -> None:
    """Validate the saved template against the fixed design contract."""

    presentation = Presentation(str(output_path))
    failures: list[str] = []
    if len(presentation.slides) != 7:
        failures.append(f"expected 7 slides, found {len(presentation.slides)}")
    if presentation.slide_width != Inches(SLIDE_WIDTH_IN):
        failures.append("slide width is not 10 inches")
    if presentation.slide_height != Inches(SLIDE_HEIGHT_IN):
        failures.append("slide height is not 5.625 inches")

    for slide_index, slide in enumerate(presentation.slides, start=1):
        picture_count = sum(
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes
        )
        if picture_count < 1:
            failures.append(f"slide {slide_index} is missing the brand logo")
        if slide_index > 1 and f"{slide_index:02d}" not in slide_text(slide):
            failures.append(f"slide {slide_index} is missing its page number")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name != FONT_NAME:
                        failures.append(
                            f"slide {slide_index} contains non-Arial text: {run.text!r}"
                        )
                    if run.font.size is None:
                        failures.append(
                            f"slide {slide_index} has an inherited font size: {run.text!r}"
                        )
                        continue
                    font_size = round(run.font.size.pt)
                    if font_size not in ALLOWED_FONT_SIZES:
                        failures.append(
                            f"slide {slide_index} uses {font_size} pt: {run.text!r}"
                        )

    evidence_slide = presentation.slides[5]
    if not any(shape.has_chart for shape in evidence_slide.shapes):
        failures.append("slide 6 is missing its native editable chart")
    if "Illustrative data" not in slide_text(evidence_slide):
        failures.append("slide 6 is missing the illustrative-data notice")

    if failures:
        details = "\n".join(f"- {failure}" for failure in failures)
        raise TemplateValidationError(f"Template validation failed:\n{details}")


def build_template(output_path: Path, *, force: bool) -> None:
    """Generate, save, and validate the canonical template."""

    resolved_output = output_path.expanduser().resolve()
    if resolved_output.exists() and not force:
        raise FileExistsError(
            f"Refusing to overwrite {resolved_output}; pass --force to replace it"
        )
    resolved_output.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_presentation(DEFAULT_LOGO)
    presentation.save(str(resolved_output))
    validate_template(resolved_output)


def parse_args() -> CliArgs:
    """Parse command-line arguments into a typed model."""

    parser = argparse.ArgumentParser(
        description="Build the canonical editable vLLM-Omni PPTX template."
    )
    parser.add_argument(
        "--output",
        type=Path,
        required=True,
        help="Destination .pptx path",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Replace an existing output file",
    )
    namespace = parser.parse_args()
    return CliArgs(output=namespace.output, force=bool(namespace.force))


def main() -> int:
    """Run the template generator."""

    args = parse_args()
    build_template(args.output, force=args.force)
    print(f"Created and validated {args.output.expanduser().resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
