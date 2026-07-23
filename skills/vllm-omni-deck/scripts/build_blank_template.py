#!/usr/bin/env python3
"""Build the canonical blank vLLM-Omni PowerPoint template."""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from build_template import (
    ALLOWED_FONT_SIZES,
    COLORS,
    FONT_NAME,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    Box,
    add_brand_plate,
    add_chevron,
    add_content_footer,
    add_content_header,
    add_dark_page_number,
    add_panel,
    add_roadmap_card,
    add_rule,
    add_text,
    set_background,
    slide_text,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Inches


SKILL_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_LOGO = SKILL_ROOT / "assets" / "vllm-omni-logo.png"


class BlankTemplateValidationError(ValueError):
    """Raised when the blank template violates its design contract."""


@dataclass(frozen=True)
class CliArgs:
    """Command-line arguments."""

    output: Path
    force: bool


def add_figure_slot(
    slide: Slide,
    box: Box,
    label: str,
    instruction: str,
) -> None:
    """Add an editable frame for an original source figure."""

    add_panel(
        slide,
        box,
        COLORS.panel,
        line_color=COLORS.rule,
        rounded=False,
    )
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.18, box.width - 0.40, 0.18),
        label,
        12,
        COLORS.blue,
        bold=True,
    )
    add_text(
        slide,
        Box(box.x + 0.35, box.y + 0.75, box.width - 0.70, 0.70),
        instruction,
        18,
        COLORS.muted,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_blank_card(
    slide: Slide,
    box: Box,
    label: str,
    accent: RGBColor,
) -> None:
    """Add one generic agenda or summary card."""

    add_panel(slide, box, COLORS.panel, line_color=COLORS.rule)
    add_rule(slide, Box(box.x, box.y, 0.07, box.height), accent)
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.15, box.width - 0.34, 0.18),
        label,
        12,
        accent,
        bold=True,
    )
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.44, box.width - 0.34, 0.30),
        "[Item heading]",
        18,
        COLORS.ink,
        bold=True,
    )
    add_text(
        slide,
        Box(box.x + 0.20, box.y + 0.86, box.width - 0.34, box.height - 1.00),
        "[One supporting line]",
        12,
        COLORS.muted,
    )


def add_cover(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 1: cover or closing."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.navy)
    add_brand_plate(slide, logo_path)
    add_text(
        slide,
        Box(0.56, 1.22, 5.62, 0.74),
        "[Presentation title]",
        36,
        COLORS.white,
        bold=True,
    )
    add_text(
        slide,
        Box(0.58, 2.35, 5.45, 0.68),
        "[One-sentence subtitle or closing statement]",
        18,
        COLORS.white,
    )
    add_rule(slide, Box(0.58, 3.30, 0.72, 0.07), COLORS.orange)
    add_text(
        slide,
        Box(0.58, 3.58, 5.20, 0.24),
        "[Presenter / team  ·  Context  ·  Month YYYY]",
        12,
        COLORS.white,
    )

    add_panel(
        slide,
        Box(6.38, 1.36, 3.06, 2.48),
        COLORS.navy,
        line_color=COLORS.rule,
    )
    add_text(
        slide,
        Box(6.62, 1.64, 2.58, 0.18),
        "OPTIONAL VISUAL",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(6.72, 2.16, 2.38, 0.78),
        "[Insert source figure unchanged or editable motif]",
        18,
        COLORS.white,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_rule(slide, Box(0.58, 5.04, 8.86, 0.012), COLORS.blue)
    add_text(
        slide,
        Box(0.58, 5.16, 7.20, 0.20),
        "BLANK vLLM-OMNI LAYOUT  ·  DELETE PLACEHOLDERS",
        12,
        COLORS.white,
        bold=True,
    )


def add_section(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 2: section or key message."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.navy)
    add_brand_plate(slide, logo_path)
    add_text(
        slide,
        Box(0.58, 1.32, 4.90, 0.20),
        "[SECTION LABEL]",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(0.58, 1.72, 5.10, 1.22),
        "[One key message]",
        36,
        COLORS.white,
        bold=True,
    )
    add_text(
        slide,
        Box(0.60, 3.20, 4.85, 0.74),
        "[One supporting sentence]",
        18,
        COLORS.white,
    )
    add_panel(
        slide,
        Box(6.15, 1.60, 3.15, 2.42),
        COLORS.navy,
        line_color=COLORS.rule,
    )
    add_text(
        slide,
        Box(6.42, 1.88, 2.61, 0.18),
        "OPTIONAL FIGURE",
        12,
        COLORS.orange,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Box(6.54, 2.39, 2.37, 0.72),
        "[Insert source figure unchanged or semantic visual]",
        18,
        COLORS.white,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_dark_page_number(slide, 2)


def add_agenda(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 3: agenda or summary."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "BLANK 03  ·  AGENDA / SUMMARY",
        "[Agenda or summary title]",
    )
    cards = (
        (Box(0.55, 1.24, 4.28, 1.52), "[01]", COLORS.blue),
        (Box(5.07, 1.24, 4.38, 1.52), "[02]", COLORS.cyan),
        (Box(0.55, 3.00, 4.28, 1.52), "[03]", COLORS.orange),
        (Box(5.07, 3.00, 4.38, 1.52), "[04]", COLORS.purple),
    )
    for box, label, accent in cards:
        add_blank_card(slide, box, label, accent)
    add_content_footer(
        slide,
        logo_path,
        3,
        "[Source / context if required]",
    )


def add_general_content(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 4: general content."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "BLANK 04  ·  GENERAL CONTENT",
        "[Assertion title]",
    )
    add_panel(slide, Box(0.55, 1.23, 3.12, 3.58), COLORS.panel, line_color=COLORS.rule)
    add_text(
        slide,
        Box(0.78, 1.49, 2.66, 0.18),
        "KEY TAKEAWAY",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(0.78, 1.85, 2.56, 0.78),
        "[One-sentence takeaway]",
        18,
        COLORS.ink,
        bold=True,
    )
    points = (
        ("01", "[Supporting point]", COLORS.blue),
        ("02", "[Supporting point]", COLORS.cyan),
        ("03", "[Supporting point]", COLORS.purple),
    )
    for index, (number, label, accent) in enumerate(points):
        y = 3.00 + index * 0.51
        add_panel(slide, Box(0.78, y, 0.42, 0.34), accent)
        add_text(
            slide,
            Box(0.78, y + 0.08, 0.42, 0.18),
            number,
            12,
            COLORS.white,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            Box(1.36, y + 0.05, 1.96, 0.22),
            label,
            12,
            COLORS.ink,
            bold=True,
        )

    add_figure_slot(
        slide,
        Box(3.98, 1.23, 5.47, 3.58),
        "FIGURE / DIAGRAM / CONTENT",
        "[Insert original source figure unchanged]",
    )
    add_text(
        slide,
        Box(4.26, 4.42, 4.91, 0.20),
        "[Caption or explanatory note]",
        12,
        COLORS.muted,
        alignment=PP_ALIGN.CENTER,
    )
    add_content_footer(
        slide,
        logo_path,
        4,
        "[Source URL]",
    )


def add_architecture(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 5: architecture or technical diagram."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "BLANK 05  ·  ARCHITECTURE / DIAGRAM",
        "[Architecture or process title]",
    )
    add_text(
        slide,
        Box(0.55, 1.12, 8.90, 0.30),
        "[One-sentence takeaway]",
        18,
        COLORS.ink,
    )

    accents = (
        COLORS.blue,
        COLORS.cyan,
        COLORS.orange,
        COLORS.cyan,
        COLORS.purple,
    )
    stage_width = 1.52
    start_x = 0.55
    gap = 0.44
    for index, accent in enumerate(accents):
        x = start_x + index * (stage_width + gap)
        add_panel(
            slide,
            Box(x, 1.72, stage_width, 1.42),
            COLORS.panel,
            line_color=accent,
        )
        add_rule(slide, Box(x, 1.72, stage_width, 0.09), accent)
        add_text(
            slide,
            Box(x + 0.13, 1.98, stage_width - 0.26, 0.18),
            f"[STAGE {index + 1}]",
            12,
            accent,
            bold=True,
        )
        add_text(
            slide,
            Box(x + 0.13, 2.31, stage_width - 0.26, 0.34),
            "[Node]",
            18,
            COLORS.ink,
            bold=True,
        )
        add_text(
            slide,
            Box(x + 0.13, 2.79, stage_width - 0.26, 0.20),
            "[Role / output]",
            12,
            COLORS.muted,
        )
        if index < len(accents) - 1:
            add_chevron(
                slide,
                Box(x + stage_width + 0.10, 2.30, 0.22, 0.30),
                COLORS.rule,
            )

    add_panel(slide, Box(0.55, 3.54, 8.90, 1.10), COLORS.navy)
    add_text(
        slide,
        Box(0.79, 3.78, 1.84, 0.18),
        "SUPPORTING BAND",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(0.79, 4.12, 7.98, 0.26),
        "[Cross-cutting concern, boundary, or annotation]",
        18,
        COLORS.white,
        bold=True,
    )
    add_content_footer(
        slide,
        logo_path,
        5,
        "[Source URL]",
    )


def add_evidence(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 6: evidence or source figure."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "BLANK 06  ·  EVIDENCE / SOURCE FIGURE",
        "[Evidence claim or figure title]",
    )
    add_text(
        slide,
        Box(0.55, 1.18, 5.40, 0.18),
        "[METRIC / UNIT / COMPARISON DIRECTION]",
        12,
        COLORS.muted,
        bold=True,
    )
    add_figure_slot(
        slide,
        Box(0.55, 1.55, 5.88, 2.94),
        "SOURCE FIGURE / CHART / TABLE",
        "[Insert original figure unchanged — scale proportionally]",
    )
    add_panel(slide, Box(0.55, 4.66, 5.88, 0.28), COLORS.orange)
    add_text(
        slide,
        Box(0.67, 4.72, 5.64, 0.18),
        "[Caption  ·  Source URL]",
        12,
        COLORS.white,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    conclusions = (
        (
            Box(6.70, 1.55, 2.75, 1.38),
            "[Conclusion 1]",
            "[What the audience should notice]",
        ),
        (
            Box(6.70, 3.11, 2.75, 1.38),
            "[Conclusion 2]",
            "[Why the evidence matters]",
        ),
    )
    for box, title, body in conclusions:
        add_panel(slide, box, COLORS.panel, line_color=COLORS.rule)
        add_text(
            slide,
            Box(box.x + 0.22, box.y + 0.22, box.width - 0.44, 0.30),
            title,
            18,
            COLORS.ink,
            bold=True,
        )
        add_text(
            slide,
            Box(box.x + 0.22, box.y + 0.70, box.width - 0.44, 0.42),
            body,
            12,
            COLORS.muted,
        )
    add_content_footer(
        slide,
        logo_path,
        6,
        "[Source URL]",
    )


def add_roadmap(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 7: roadmap or status."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "BLANK 07  ·  ROADMAP / STATUS",
        "[Roadmap or status title]",
    )
    add_text(
        slide,
        Box(0.55, 1.12, 8.90, 0.28),
        "[One-sentence takeaway]",
        18,
        COLORS.ink,
    )
    cards = (
        (
            Box(0.55, 1.63, 2.80, 2.65),
            "[PHASE 1]",
            "[Phase title]",
            ("[Item]", "[Item]"),
            COLORS.blue,
        ),
        (
            Box(3.60, 1.63, 2.80, 2.65),
            "[PHASE 2]",
            "[Phase title]",
            ("[Item]", "[Item]"),
            COLORS.orange,
        ),
        (
            Box(6.65, 1.63, 2.80, 2.65),
            "[PHASE 3]",
            "[Phase title]",
            ("[Item]", "[Item]"),
            COLORS.purple,
        ),
    )
    for box, label, title, items, accent in cards:
        add_roadmap_card(slide, box, label, title, items, accent)
    add_panel(slide, Box(0.55, 4.51, 8.90, 0.46), COLORS.navy)
    add_text(
        slide,
        Box(0.76, 4.64, 1.45, 0.18),
        "DECISION / RISK",
        12,
        COLORS.orange,
        bold=True,
    )
    add_text(
        slide,
        Box(2.35, 4.60, 6.72, 0.24),
        "[Decision, dependency, or risk statement]",
        18,
        COLORS.white,
        bold=True,
    )
    add_content_footer(
        slide,
        logo_path,
        7,
        "[Source / owner / date]",
    )


def add_white_canvas(
    presentation: PresentationType,
    logo_path: Path,
) -> None:
    """Build blank layout 8: a minimally structured white canvas."""

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide, COLORS.white)
    add_content_header(
        slide,
        "BLANK 08  ·  WHITE CANVAS",
        "[Slide title]",
    )
    add_content_footer(
        slide,
        logo_path,
        8,
        "[Source / context if required]",
    )


def build_presentation(logo_path: Path) -> PresentationType:
    """Build the complete blank eight-slide presentation."""

    if not logo_path.is_file():
        raise FileNotFoundError(f"Missing vLLM-Omni logo: {logo_path}")
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    presentation.core_properties.title = "vLLM-Omni blank presentation template"
    presentation.core_properties.subject = "Eight reusable vLLM-Omni layouts"
    presentation.core_properties.author = "vLLM-Omni"
    add_cover(presentation, logo_path)
    add_section(presentation, logo_path)
    add_agenda(presentation, logo_path)
    add_general_content(presentation, logo_path)
    add_architecture(presentation, logo_path)
    add_evidence(presentation, logo_path)
    add_roadmap(presentation, logo_path)
    add_white_canvas(presentation, logo_path)
    return presentation


def validate_blank_template(output_path: Path) -> None:
    """Validate the blank template against the fixed design contract."""

    presentation = Presentation(str(output_path))
    failures: list[str] = []
    if len(presentation.slides) != 8:
        failures.append(f"expected 8 slides, found {len(presentation.slides)}")
    if presentation.slide_width != Inches(SLIDE_WIDTH_IN):
        failures.append("slide width is not 10 inches")
    if presentation.slide_height != Inches(SLIDE_HEIGHT_IN):
        failures.append("slide height is not 5.625 inches")

    for slide_index, slide in enumerate(presentation.slides, start=1):
        text = slide_text(slide)
        picture_count = sum(
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes
        )
        if picture_count < 1:
            failures.append(f"slide {slide_index} is missing the brand logo")
        if "[" not in text or "]" not in text:
            failures.append(f"slide {slide_index} is missing editable placeholders")
        if slide_index > 1 and f"{slide_index:02d}" not in text:
            failures.append(f"slide {slide_index} is missing its page number")
        if "adapted from" in text.lower():
            failures.append(f"slide {slide_index} permits source-figure adaptation")
        for shape in slide.shapes:
            if shape.has_chart:
                failures.append(
                    f"slide {slide_index} contains example chart data in blank mode"
                )
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name != FONT_NAME:
                        failures.append(
                            f"slide {slide_index} contains non-Arial text: "
                            f"{run.text!r}"
                        )
                    if run.font.size is None:
                        failures.append(
                            f"slide {slide_index} has inherited font size: "
                            f"{run.text!r}"
                        )
                        continue
                    font_size = round(run.font.size.pt)
                    if font_size not in ALLOWED_FONT_SIZES:
                        failures.append(
                            f"slide {slide_index} uses {font_size} pt: "
                            f"{run.text!r}"
                        )

    if "Insert original source figure unchanged" not in slide_text(
        presentation.slides[3]
    ):
        failures.append("slide 4 is missing its intact source-figure slot")
    if "Insert original figure unchanged" not in slide_text(presentation.slides[5]):
        failures.append("slide 6 is missing its intact evidence-figure slot")
    white_canvas = presentation.slides[7]
    white_canvas_text = slide_text(white_canvas)
    required_white_canvas_text = (
        "BLANK 08  ·  WHITE CANVAS",
        "[Slide title]",
        "[Source / context if required]",
    )
    for required_text in required_white_canvas_text:
        if required_text not in white_canvas_text:
            failures.append(f"slide 8 is missing {required_text!r}")
    if white_canvas_text.count("08") < 2:
        failures.append("slide 8 is missing its page number")
    if len(white_canvas.shapes) != 7:
        failures.append("slide 8 body contains unexpected objects")
    if white_canvas.background.fill.fore_color.rgb != COLORS.white:
        failures.append("slide 8 background is not white")

    if failures:
        details = "\n".join(f"- {failure}" for failure in failures)
        raise BlankTemplateValidationError(
            f"Blank template validation failed:\n{details}"
        )


def build_blank_template(output_path: Path, *, force: bool) -> None:
    """Generate, save, and validate the blank template."""

    resolved_output = output_path.expanduser().resolve()
    if resolved_output.exists() and not force:
        raise FileExistsError(
            f"Refusing to overwrite {resolved_output}; pass --force to replace it"
        )
    resolved_output.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_presentation(DEFAULT_LOGO)
    presentation.save(str(resolved_output))
    validate_blank_template(resolved_output)


def parse_args() -> CliArgs:
    """Parse command-line arguments into a typed model."""

    parser = argparse.ArgumentParser(
        description="Build the blank vLLM-Omni PPTX template."
    )
    parser.add_argument(
        "--output",
        type=Path,
        required=True,
        help="Destination .pptx path",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Replace an existing output file",
    )
    namespace = parser.parse_args()
    return CliArgs(output=namespace.output, force=bool(namespace.force))


def main() -> int:
    """Run the blank-template generator."""

    args = parse_args()
    build_blank_template(args.output, force=args.force)
    print(f"Created and validated {args.output.expanduser().resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
