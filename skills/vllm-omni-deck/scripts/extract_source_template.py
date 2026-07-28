#!/usr/bin/env python3
"""Extract a compact editable template from the public vLLM-Omni deck."""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.text.text import _Paragraph
from pptx.util import Inches, Pt

FONT_NAME = "Arial"
SLIDE_WIDTH_IN = 10.0
SLIDE_HEIGHT_IN = 5.625
MIN_FONT_SIZE_PT = 12

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x78, 0x90, 0x9C)
BLUE = RGBColor(0x42, 0x85, 0xF4)
ORANGE = RGBColor(0xFF, 0xAB, 0x40)
TEAL = RGBColor(0x00, 0x97, 0xA7)
YELLOW = RGBColor(0xEE, 0xFF, 0x41)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)

SOURCE_SIGNATURES = {
    1: "meetup",
    9: "Goal of the walkthrough",
    34: "Summary",
    42: "Diffusion models",
    45: "vLLM Networking Hour!",
}
OUTPUT_ORDER = (1, 34, 9, 42, 45)
OUTPUT_TITLES = (
    "[Presentation title]",
    "Contents",
    "[Slide title]",
    "Chart color guidance",
    "Thank you.",
)
RELATIONSHIP_NAMESPACE = (
    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
)


class TemplateExtractionError(ValueError):
    """Raised when the source or extracted template violates the contract."""


@dataclass(frozen=True)
class CliArgs:
    """Command-line arguments."""

    source: Path
    output: Path
    force: bool


def parse_args() -> CliArgs:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--force", action="store_true")
    args = parser.parse_args()
    return CliArgs(args.source, args.output, args.force)


def slide_title(slide: Slide) -> str:
    """Return the visible slide title."""

    title = slide.shapes.title
    return title.text.strip() if title is not None else ""


def validate_source(presentation: PresentationType) -> None:
    """Confirm the input matches the expected public source deck."""

    width = presentation.slide_width / 914400
    height = presentation.slide_height / 914400
    if abs(width - SLIDE_WIDTH_IN) > 0.01 or abs(height - SLIDE_HEIGHT_IN) > 0.01:
        raise TemplateExtractionError(
            f"Expected a {SLIDE_WIDTH_IN} x {SLIDE_HEIGHT_IN} inch source, "
            f"got {width:.3f} x {height:.3f}"
        )
    if len(presentation.slides) < max(SOURCE_SIGNATURES):
        raise TemplateExtractionError(
            f"Source has {len(presentation.slides)} slides; "
            f"expected at least {max(SOURCE_SIGNATURES)}"
        )
    for slide_number, expected_title in SOURCE_SIGNATURES.items():
        actual_title = slide_title(presentation.slides[slide_number - 1])
        if actual_title != expected_title:
            raise TemplateExtractionError(
                f"Slide {slide_number} title changed: "
                f"expected {expected_title!r}, got {actual_title!r}"
            )


def delete_slide(presentation: PresentationType, index: int) -> None:
    """Delete one slide and its package relationship."""

    slide_id = presentation.slides._sldIdLst[index]
    presentation.part.drop_rel(slide_id.rId)
    del presentation.slides._sldIdLst[index]


def keep_and_reorder_slides(presentation: PresentationType) -> None:
    """Keep the five source exemplars and put them in template order."""

    selected_ids = {
        number: presentation.slides[number - 1].slide_id for number in OUTPUT_ORDER
    }
    keep_numbers = set(OUTPUT_ORDER)
    for index in reversed(range(len(presentation.slides))):
        if index + 1 not in keep_numbers:
            delete_slide(presentation, index)

    slide_id_list = presentation.slides._sldIdLst
    by_id = {int(slide_id.id): slide_id for slide_id in slide_id_list}
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for source_number in OUTPUT_ORDER:
        slide_id_list.append(by_id[selected_ids[source_number]])


def related_relationship_ids(shape: BaseShape) -> set[str]:
    """Collect relationships referenced by one shape subtree."""

    relationship_ids: set[str] = set()
    for element in shape._element.iter():
        for attribute, value in element.attrib.items():
            if attribute.startswith(RELATIONSHIP_NAMESPACE):
                relationship_ids.add(value)
    return relationship_ids


def remove_shape(slide: Slide, shape: BaseShape) -> None:
    """Remove a shape and any relationships used only by that shape."""

    relationship_ids = related_relationship_ids(shape)
    shape._element.getparent().remove(shape._element)
    for relationship_id in relationship_ids:
        if relationship_id in slide.part.rels:
            slide.part.drop_rel(relationship_id)


def remove_bullet_markup(paragraph: _Paragraph) -> None:
    """Remove local bullet or numbering markup."""

    properties = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont", "a:buClr"):
        for element in properties.findall(qn(tag)):
            properties.remove(element)


def clear_bullet(paragraph: _Paragraph) -> None:
    """Explicitly disable inherited bullet or numbering markup."""

    remove_bullet_markup(paragraph)
    paragraph._p.get_or_add_pPr().append(OxmlElement("a:buNone"))


def set_bullet(paragraph: _Paragraph, *, numbered: bool = False) -> None:
    """Apply a native bullet or automatic number."""

    remove_bullet_markup(paragraph)
    properties = paragraph._p.get_or_add_pPr()
    properties.set("marL", str(int(Inches(0.42))))
    properties.set("indent", str(-int(Inches(0.20))))
    marker = OxmlElement("a:buAutoNum" if numbered else "a:buChar")
    if numbered:
        marker.set("type", "arabicPeriod")
    else:
        marker.set("char", "\u2022")
    properties.append(marker)


def style_paragraph(
    paragraph: _Paragraph,
    text: str,
    font_size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    numbered: bool = False,
    bullet: bool = False,
    space_after: int = 8,
) -> None:
    """Populate and style one paragraph."""

    paragraph.alignment = alignment
    paragraph.level = 0
    paragraph.space_after = Pt(space_after)
    paragraph.line_spacing = 1.1
    if numbered or bullet:
        set_bullet(paragraph, numbered=numbered)
    else:
        clear_bullet(paragraph)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def reset_text_frame(shape: BaseShape) -> None:
    """Clear a text frame and normalize its box behavior."""

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)


def add_text_box(
    slide: Slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> BaseShape:
    """Add one editable text box."""

    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    reset_text_frame(shape)
    shape.text_frame.vertical_anchor = vertical_anchor
    style_paragraph(
        shape.text_frame.paragraphs[0],
        text,
        font_size,
        color,
        bold=bold,
        alignment=alignment,
        space_after=0,
    )
    return shape


def add_rule(
    slide: Slide,
    left: float,
    top: float,
    width: float,
    height: float,
    color: RGBColor,
) -> BaseShape:
    """Add a flat rectangular rule."""

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_white_background(slide: Slide) -> None:
    """Set a plain white slide background."""

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def body_placeholder(slide: Slide) -> BaseShape:
    """Return the body placeholder on a source slide."""

    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            return shape
    raise TemplateExtractionError("Slide does not contain a body placeholder")


def move_slide_number(slide: Slide) -> None:
    """Place the native slide-number field immediately left of the footer logo."""

    for shape in slide.placeholders:
        if shape.placeholder_format.type != PP_PLACEHOLDER.SLIDE_NUMBER:
            continue
        shape.left = Inches(6.78)
        shape.top = Inches(5.16)
        shape.width = Inches(0.22)
        shape.height = Inches(0.24)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def build_cover(slide: Slide) -> None:
    """Convert the public-deck cover into an editable template cover."""

    set_white_background(slide)
    title = slide.shapes.title
    if title is None:
        raise TemplateExtractionError("Cover slide has no title placeholder")
    title.left = Inches(1.30)
    title.top = Inches(1.96)
    title.width = Inches(7.80)
    title.height = Inches(0.86)
    reset_text_frame(title)
    title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    style_paragraph(
        title.text_frame.paragraphs[0],
        OUTPUT_TITLES[0],
        36,
        INK,
        bold=True,
        space_after=0,
    )

    subtitle = next(
        (
            shape
            for shape in slide.placeholders
            if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
        ),
        None,
    )
    if subtitle is None:
        raise TemplateExtractionError("Cover slide has no subtitle placeholder")
    subtitle.left = Inches(1.30)
    subtitle.top = Inches(3.02)
    subtitle.width = Inches(6.90)
    subtitle.height = Inches(0.90)
    reset_text_frame(subtitle)
    style_paragraph(
        subtitle.text_frame.paragraphs[0],
        "[Presenter / team]",
        18,
        INK,
        space_after=6,
    )
    style_paragraph(
        subtitle.text_frame.add_paragraph(),
        "[Month YYYY]",
        12,
        MUTED,
        space_after=0,
    )


def build_contents(slide: Slide) -> None:
    """Convert the summary exemplar into a sparse contents page."""

    set_white_background(slide)
    move_slide_number(slide)
    title = slide.shapes.title
    title.text = ""
    reset_text_frame(title)
    style_paragraph(
        title.text_frame.paragraphs[0],
        OUTPUT_TITLES[1],
        28,
        INK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        space_after=0,
    )

    body = body_placeholder(slide)
    body.left = Inches(1.22)
    body.top = Inches(1.30)
    body.width = Inches(7.60)
    body.height = Inches(3.42)
    reset_text_frame(body)
    for index in range(4):
        paragraph = (
            body.text_frame.paragraphs[0]
            if index == 0
            else body.text_frame.add_paragraph()
        )
        style_paragraph(
            paragraph,
            " [Section title]",
            18,
            INK,
            numbered=True,
            space_after=18,
        )


def build_body(slide: Slide) -> None:
    """Convert the walkthrough-goal exemplar into a reusable body slide."""

    set_white_background(slide)
    move_slide_number(slide)
    title = slide.shapes.title
    reset_text_frame(title)
    style_paragraph(
        title.text_frame.paragraphs[0],
        OUTPUT_TITLES[2],
        28,
        INK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        space_after=0,
    )

    body = body_placeholder(slide)
    body.left = Inches(0.78)
    body.top = Inches(1.28)
    body.width = Inches(8.44)
    body.height = Inches(3.34)
    reset_text_frame(body)
    style_paragraph(
        body.text_frame.paragraphs[0],
        "[Key takeaway]",
        18,
        BLUE,
        bold=True,
        space_after=18,
    )
    for _ in range(3):
        style_paragraph(
            body.text_frame.add_paragraph(),
            "[Supporting point]",
            18,
            INK,
            bullet=True,
            space_after=12,
        )
    add_text_box(
        slide,
        0.36,
        4.83,
        6.45,
        0.24,
        "[Source / context]",
        12,
        MUTED,
    )


def add_comparison_chart(slide: Slide) -> BaseShape:
    """Add a self-contained native chart with neutral illustrative data."""

    data = ChartData()
    data.categories = ("Option A", "Option B", "Option C")
    data.add_series("Baseline", (64, 78, 72))
    data.add_series("Improved", (82, 91, 87))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.55),
        Inches(1.50),
        Inches(4.08),
        Inches(3.22),
        data,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.major_unit = 20
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.name = FONT_NAME
    chart.value_axis.tick_labels.font.size = Pt(12)
    chart.category_axis.tick_labels.font.name = FONT_NAME
    chart.category_axis.tick_labels.font.size = Pt(12)
    for series, color in zip(chart.series, (BLUE, ORANGE), strict=True):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.format.line.color.rgb = color
    return chart_shape


def set_cell_text(cell, text: str, *, bold: bool = False, color: RGBColor = INK) -> None:
    """Set editable table-cell text."""

    frame = cell.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    style_paragraph(
        frame.paragraphs[0],
        text,
        12,
        color,
        bold=bold,
        alignment=PP_ALIGN.LEFT,
        space_after=0,
    )


def add_palette_table(slide: Slide) -> None:
    """Add an editable palette table beside the native chart."""

    rows = (
        ("Primary", BLUE, "#4285F4", "Main series"),
        ("Ink", INK, "#212121", "Text / baseline"),
        ("Slate", MUTED, "#78909C", "Secondary"),
        ("Orange", ORANGE, "#FFAB40", "Highlight"),
        ("Teal", TEAL, "#0097A7", "Supporting"),
        ("Yellow", YELLOW, "#EEFF41", "Sparse accent"),
    )
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        4,
        Inches(5.05),
        Inches(1.50),
        Inches(4.40),
        Inches(3.22),
    )
    table = table_shape.table
    widths = (1.05, 0.72, 1.10, 1.53)
    for column, width in zip(table.columns, widths, strict=True):
        column.width = Inches(width)

    headers = ("Role", "Color", "Hex", "Use")
    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        set_cell_text(cell, header, bold=True, color=WHITE)

    for row_index, (role, swatch, hex_value, use) in enumerate(rows, start=1):
        values = (role, "", hex_value, use)
        for column_index, value in enumerate(values):
            cell = table.cell(row_index, column_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else LIGHT_GRAY
            set_cell_text(cell, value)
        swatch_cell = table.cell(row_index, 1)
        swatch_cell.fill.solid()
        swatch_cell.fill.fore_color.rgb = swatch


def build_chart_guidance(slide: Slide) -> None:
    """Convert the source chart slide into chart and palette guidance."""

    set_white_background(slide)
    move_slide_number(slide)
    title = slide.shapes.title
    reset_text_frame(title)
    style_paragraph(
        title.text_frame.paragraphs[0],
        OUTPUT_TITLES[3],
        28,
        INK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        space_after=0,
    )

    chart_shapes = [
        shape for shape in slide.shapes if getattr(shape, "has_chart", False)
    ]
    if len(chart_shapes) != 2:
        raise TemplateExtractionError(
            f"Expected two source charts on slide 42, found {len(chart_shapes)}"
        )
    for chart_shape in chart_shapes:
        remove_shape(slide, chart_shape)
    add_comparison_chart(slide)

    add_text_box(
        slide,
        0.58,
        1.15,
        4.00,
        0.26,
        "Illustrative comparison",
        18,
        INK,
        bold=True,
    )
    add_text_box(
        slide,
        5.08,
        1.15,
        4.30,
        0.26,
        "Source palette",
        18,
        INK,
        bold=True,
    )
    add_palette_table(slide)
    add_text_box(
        slide,
        0.36,
        4.83,
        7.30,
        0.24,
        "Illustrative data | Replace labels and values before use",
        12,
        MUTED,
    )


def build_closing(slide: Slide) -> None:
    """Convert the networking slide into a clean closing page."""

    title = slide.shapes.title
    for shape in list(slide.shapes):
        if shape._element is title._element:
            continue
        if (
            shape.is_placeholder
            and shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER
        ):
            continue
        remove_shape(slide, shape)

    set_white_background(slide)
    title.left = Inches(0.72)
    title.top = Inches(1.30)
    title.width = Inches(4.25)
    title.height = Inches(0.86)
    reset_text_frame(title)
    title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    style_paragraph(
        title.text_frame.paragraphs[0],
        OUTPUT_TITLES[4],
        36,
        INK,
        bold=True,
        space_after=0,
    )
    add_rule(slide, 0.73, 2.43, 0.06, 1.33, BLUE)
    add_text_box(
        slide,
        1.02,
        2.47,
        3.95,
        1.22,
        "[Closing statement]",
        18,
        INK,
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide,
        5.55,
        1.55,
        3.70,
        0.36,
        "[Project / team]",
        18,
        INK,
        bold=True,
    )
    contact = add_text_box(
        slide,
        5.55,
        2.10,
        3.70,
        1.68,
        "[Contact name]",
        12,
        MUTED,
    )
    for text in (
        "[Email or messaging handle]",
        "[Repository or documentation URL]",
    ):
        style_paragraph(
            contact.text_frame.add_paragraph(),
            text,
            12,
            MUTED,
            space_after=10,
        )


def extract_template(source: Path, output: Path) -> None:
    """Create the five-slide source-derived template."""

    presentation = Presentation(source)
    validate_source(presentation)
    keep_and_reorder_slides(presentation)

    builders = (
        build_cover,
        build_contents,
        build_body,
        build_chart_guidance,
        build_closing,
    )
    for slide, builder in zip(presentation.slides, builders, strict=True):
        builder(slide)

    presentation.core_properties.title = "vLLM-Omni source-derived template"
    presentation.core_properties.subject = (
        "Cover, contents, body, chart guidance, and closing template"
    )
    presentation.core_properties.comments = (
        "Extracted from vLLM-Omni Slides (Public) 2026-04 latest.pptx"
    )
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def slide_text(slide: Slide) -> str:
    """Collect slide-local text for validation."""

    return "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def validate_output(path: Path) -> None:
    """Validate the extracted template structure and editable content."""

    presentation = Presentation(path)
    if len(presentation.slides) != 5:
        raise TemplateExtractionError(
            f"Expected five output slides, found {len(presentation.slides)}"
        )
    actual_titles = tuple(slide_title(slide) for slide in presentation.slides)
    if actual_titles != OUTPUT_TITLES:
        raise TemplateExtractionError(
            f"Unexpected output titles: {actual_titles!r}"
        )

    chart_slide = presentation.slides[3]
    charts = [
        shape for shape in chart_slide.shapes if getattr(shape, "has_chart", False)
    ]
    tables = [
        shape for shape in chart_slide.shapes if getattr(shape, "has_table", False)
    ]
    if len(charts) != 1 or len(tables) != 1:
        raise TemplateExtractionError(
            "Chart guidance must contain one native chart and one native table"
        )

    stale_text = (
        "meetup",
        "Apr 2026",
        "Goal of the walkthrough",
        "Summary",
        "Diffusion models",
        "Networking Hour",
        "Qwen",
        "Wan2.2",
    )
    full_text = "\n".join(slide_text(slide) for slide in presentation.slides)
    for value in stale_text:
        if value in full_text:
            raise TemplateExtractionError(f"Stale source text remains: {value!r}")

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > presentation.slide_width
                or shape.top + shape.height > presentation.slide_height
            ):
                raise TemplateExtractionError(
                    f"Slide {slide_number} shape exceeds the canvas: {shape.name}"
                )
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if (
                        run.text.strip()
                        and run.font.size is not None
                        and run.font.size.pt < MIN_FONT_SIZE_PT
                    ):
                        raise TemplateExtractionError(
                            f"Slide {slide_number} contains "
                            f"{run.font.size.pt:g} pt text in {shape.name}"
                        )


def main() -> None:
    args = parse_args()
    if not args.source.is_file():
        raise SystemExit(f"Source file not found: {args.source}")
    if args.output.exists() and not args.force:
        raise SystemExit(f"Output already exists: {args.output}; pass --force")
    extract_template(args.source, args.output)
    validate_output(args.output)
    print(args.output)


if __name__ == "__main__":
    main()
